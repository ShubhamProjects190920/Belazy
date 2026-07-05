"""
Module 8 — Report Generation Service
PDF (reportlab), PowerPoint (python-pptx), Email (aiosmtplib with attachment)
"""
import html
import logging
import re
from datetime import date
from email.mime.application import MIMEApplication
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
from io import BytesIO

import aiosmtplib
from reportlab.lib import colors
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import cm
from reportlab.platypus import (
    HRFlowable,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from app.core.config import settings
from app.models.ai_report import AIReport
from app.models.project import Project

logger = logging.getLogger(__name__)

# ── Brand colours ─────────────────────────────────────────────────────────────

BRAND_BLUE  = colors.HexColor('#1e3a5f')
BRAND_LIGHT = colors.HexColor('#e8f0fe')
GREY_LIGHT  = colors.HexColor('#f8fafc')
GREY_BORDER = colors.HexColor('#d1d5db')


# ── Text helpers ──────────────────────────────────────────────────────────────

_EMOJI_MAP = {
    '🟢': 'ON TRACK', '🟡': 'AT RISK', '🔴': 'DELAYED',
    '✓': 'Yes', '✗': 'No', '⚠️': 'Warning', '✅': 'Done',
    '📋': '', '📊': '', '📄': '',
}


def _pdf_clean(text: str) -> str:
    """Strip emoji and replace ₹ (not in Helvetica) with Rs. for PDF compatibility."""
    for ch, rep in _EMOJI_MAP.items():
        text = text.replace(ch, rep)
    text = text.replace('₹', 'Rs.')
    return text


def _safe_para(text: str) -> str:
    """Escape XML entities then apply inline markdown → reportlab tags."""
    text = _pdf_clean(text)
    text = html.escape(text)
    text = re.sub(r'\*\*(.+?)\*\*', r'<b>\1</b>', text)
    text = re.sub(r'`([^`]+)`', r'<font name="Courier" size="8">\1</font>', text)
    return text


def _fmt_budget(val) -> str:
    if val is None:
        return '-'
    return f'Rs.{float(val):,.0f}'


def _is_sep_row(line: str) -> bool:
    return bool(re.match(r'^\|[\s\-:|]+\|$', line.strip()))


# ── PDF Styles ─────────────────────────────────────────────────────────────────

def _styles():
    base = getSampleStyleSheet()
    return {
        'h1': ParagraphStyle('h1', parent=base['Normal'],
                              fontSize=14, fontName='Helvetica-Bold',
                              textColor=BRAND_BLUE, spaceBefore=10, spaceAfter=4),
        'h2': ParagraphStyle('h2', parent=base['Normal'],
                              fontSize=11, fontName='Helvetica-Bold',
                              textColor=BRAND_BLUE, spaceBefore=8, spaceAfter=2),
        'body': ParagraphStyle('body', parent=base['Normal'],
                                fontSize=9, leading=14, spaceAfter=3),
        'bullet': ParagraphStyle('bullet', parent=base['Normal'],
                                  fontSize=9, leading=13, leftIndent=14, spaceAfter=2),
        'meta': ParagraphStyle('meta', parent=base['Normal'],
                                fontSize=9, fontName='Helvetica-Bold',
                                textColor=BRAND_BLUE, spaceAfter=2),
        'footer': ParagraphStyle('footer', parent=base['Normal'],
                                  fontSize=7, textColor=colors.grey, alignment=1),
        'th': ParagraphStyle('th', parent=base['Normal'],
                              fontSize=8, fontName='Helvetica-Bold', textColor=BRAND_BLUE),
        'td': ParagraphStyle('td', parent=base['Normal'], fontSize=8, leading=10),
        'kv': ParagraphStyle('kv', parent=base['Normal'],
                              fontSize=20, fontName='Helvetica-Bold',
                              textColor=BRAND_BLUE, alignment=1),
        'kl': ParagraphStyle('kl', parent=base['Normal'],
                              fontSize=8, textColor=colors.grey, alignment=1),
    }


# ── Markdown → Story parser ────────────────────────────────────────────────────

def _md_to_story(content: str, st: dict) -> list:
    story = []
    lines = content.split('\n')
    i = 0

    while i < len(lines):
        line = lines[i]

        if line.startswith('## '):
            story.append(Spacer(1, 0.1 * cm))
            story.append(HRFlowable(width='100%', thickness=1, color=BRAND_BLUE, spaceAfter=2))
            story.append(Paragraph(_safe_para(line[3:]), st['h1']))

        elif line.startswith('### '):
            story.append(Paragraph(_safe_para(line[4:]), st['h2']))

        elif line.startswith('| '):
            # Collect all table rows
            trows = []
            while i < len(lines) and lines[i].startswith('|'):
                trows.append(lines[i])
                i += 1

            tdata = []
            is_header = True
            for tr in trows:
                if _is_sep_row(tr):
                    continue
                cells = [c.strip() for c in tr.split('|') if c.strip()]
                if not cells:
                    continue
                sty = st['th'] if is_header else st['td']
                tdata.append([Paragraph(_safe_para(c), sty) for c in cells])
                is_header = False

            if tdata:
                n = max(len(r) for r in tdata)
                cw = 17.0 * cm / n
                tbl = Table(tdata, colWidths=[cw] * n, repeatRows=1)
                tbl.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), BRAND_LIGHT),
                    ('GRID', (0, 0), (-1, -1), 0.5, GREY_BORDER),
                    ('PADDING', (0, 0), (-1, -1), 5),
                    ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, GREY_LIGHT]),
                ]))
                story.append(tbl)
                story.append(Spacer(1, 0.2 * cm))
            continue

        elif line.startswith('- '):
            while i < len(lines) and lines[i].startswith('- '):
                story.append(Paragraph(f'• {_safe_para(lines[i][2:])}', st['bullet']))
                i += 1
            continue

        elif line.strip() in ('', '---'):
            story.append(Spacer(1, 0.1 * cm))

        else:
            stripped = line.strip()
            if stripped:
                if stripped.startswith('**') and ':' in stripped:
                    story.append(Paragraph(_safe_para(stripped), st['meta']))
                else:
                    story.append(Paragraph(_safe_para(stripped), st['body']))

        i += 1

    return story


# ── PDF: AI Report ─────────────────────────────────────────────────────────────

def generate_ai_report_pdf(report: AIReport, company_name: str) -> BytesIO:
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        rightMargin=2 * cm, leftMargin=2 * cm,
        topMargin=1.5 * cm, bottomMargin=2 * cm,
        title=report.title,
    )
    st = _styles()
    story = []

    # Header bar
    hdr_st = ParagraphStyle('hd', fontSize=10, fontName='Helvetica-Bold', textColor=colors.white)
    sub_st = ParagraphStyle('hs', fontSize=8, textColor=colors.white, alignment=2)
    hdr = Table(
        [[Paragraph(html.escape(_pdf_clean(company_name)), hdr_st),
          Paragraph(report.created_at.strftime('%d %b %Y, %H:%M'), sub_st)]],
        colWidths=[11 * cm, 6 * cm],
    )
    hdr.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), BRAND_BLUE),
        ('PADDING', (0, 0), (-1, -1), 8),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    story.append(hdr)
    story.append(Spacer(1, 0.5 * cm))

    story.extend(_md_to_story(report.content, st))

    story.append(Spacer(1, 0.4 * cm))
    story.append(HRFlowable(width='100%', thickness=0.5, color=colors.lightgrey))
    story.append(Paragraph(
        f'AI Project Controls Platform  |  {html.escape(_pdf_clean(company_name))}  |  Confidential',
        st['footer'],
    ))

    doc.build(story)
    buf.seek(0)
    return buf


# ── PDF: Project Summary ───────────────────────────────────────────────────────

def generate_project_summary_pdf(projects: list, company_name: str) -> BytesIO:
    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        rightMargin=2 * cm, leftMargin=2 * cm,
        topMargin=1.5 * cm, bottomMargin=2 * cm,
    )
    st = _styles()
    story = []
    today = date.today().strftime('%d %B %Y')

    # Cover header
    cmp = _pdf_clean(company_name)
    cover = Table(
        [[Paragraph(html.escape(cmp), ParagraphStyle('ch', fontSize=13, fontName='Helvetica-Bold', textColor=colors.white)),
          Paragraph(f'Project Summary Report\n{today}',
                    ParagraphStyle('cs', fontSize=9, textColor=colors.HexColor('#bfdbfe'), alignment=2))]],
        colWidths=[10 * cm, 7 * cm],
    )
    cover.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, -1), BRAND_BLUE),
        ('PADDING', (0, 0), (-1, -1), 12),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    story.append(cover)
    story.append(Spacer(1, 0.6 * cm))

    # KPIs
    total = len(projects)
    active = sum(1 for p in projects if p.status.value == 'active')
    completed = sum(1 for p in projects if p.status.value == 'completed')
    budget_sum = sum(float(p.budget) for p in projects if p.budget)

    kpi = Table(
        [
            [Paragraph('<b>Total</b>', st['kl']),
             Paragraph('<b>Active</b>', st['kl']),
             Paragraph('<b>Completed</b>', st['kl']),
             Paragraph('<b>Budget</b>', st['kl'])],
            [Paragraph(str(total), st['kv']),
             Paragraph(str(active), ParagraphStyle('kv2', parent=st['kv'], textColor=colors.HexColor('#059669'))),
             Paragraph(str(completed), ParagraphStyle('kv3', parent=st['kv'], textColor=colors.HexColor('#475569'))),
             Paragraph(_fmt_budget(budget_sum or None),
                       ParagraphStyle('kv4', parent=st['kv'], fontSize=14, textColor=colors.HexColor('#d97706')))],
        ],
        colWidths=[4.25 * cm] * 4,
    )
    kpi.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), BRAND_LIGHT),
        ('BOX', (0, 0), (-1, -1), 1, BRAND_BLUE),
        ('INNERGRID', (0, 0), (-1, -1), 0.5, GREY_BORDER),
        ('PADDING', (0, 0), (-1, -1), 8),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN', (0, 0), (-1, -1), 'MIDDLE'),
    ]))
    story.append(kpi)
    story.append(Spacer(1, 0.6 * cm))

    # Projects table
    story.append(Paragraph('All Projects', st['h1']))
    if not projects:
        story.append(Paragraph('No projects found.', st['body']))
    else:
        headers = ['Project Name', 'Client', 'Status', 'Priority', 'Budget', 'Start Date']
        rows = [[Paragraph(h, st['th']) for h in headers]]
        for p in projects:
            rows.append([
                Paragraph(html.escape(_pdf_clean(p.name)), st['td']),
                Paragraph(html.escape(_pdf_clean(p.client_name or '-')), st['td']),
                Paragraph(p.status.value.replace('_', ' ').title(), st['td']),
                Paragraph(p.priority.value.title(), st['td']),
                Paragraph(_fmt_budget(p.budget), st['td']),
                Paragraph(str(p.start_date) if p.start_date else '-', st['td']),
            ])
        tbl = Table(rows, colWidths=[4.8*cm, 3.2*cm, 2.4*cm, 2*cm, 2.6*cm, 2*cm], repeatRows=1)
        tbl.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), BRAND_LIGHT),
            ('GRID', (0, 0), (-1, -1), 0.5, GREY_BORDER),
            ('PADDING', (0, 0), (-1, -1), 5),
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, GREY_LIGHT]),
        ]))
        story.append(tbl)

    story.append(Spacer(1, 0.5 * cm))
    story.append(HRFlowable(width='100%', thickness=0.5, color=colors.lightgrey))
    story.append(Paragraph(
        f'AI Project Controls Platform  |  {html.escape(cmp)}  |  {today}  |  Confidential',
        st['footer'],
    ))

    doc.build(story)
    buf.seek(0)
    return buf


# ── PowerPoint: Project Status Deck ───────────────────────────────────────────

def generate_project_ppt(projects: list, company_name: str) -> BytesIO:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor

    BRAND = RGBColor(0x1e, 0x3a, 0x5f)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    LIGHT = RGBColor(0xe8, 0xf0, 0xfe)
    SLATE = RGBColor(0x64, 0x74, 0x8b)
    SKY   = RGBColor(0xba, 0xe6, 0xfd)
    GREEN = RGBColor(0x05, 0x96, 0x69)
    AMBER = RGBColor(0xd9, 0x77, 0x06)

    prs = Presentation()
    prs.slide_width  = Emu(9_144_000)   # 16:9
    prs.slide_height = Emu(5_143_500)
    blank = prs.slide_layouts[6]

    def slide():
        return prs.slides.add_slide(blank)

    def bg(s, rgb):
        f = s.background.fill
        f.solid()
        f.fore_color.rgb = rgb

    def tb(s, text, l, t, w, h, size=14, bold=False, color=WHITE, align='left'):
        from pptx.enum.text import PP_ALIGN
        bx = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = bx.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        p.alignment = {'left': PP_ALIGN.LEFT, 'center': PP_ALIGN.CENTER, 'right': PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)

    def rect(s, l, t, w, h, fill_rgb, border_rgb=None):
        shape = s.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
        if border_rgb:
            shape.line.color.rgb = border_rgb
        else:
            shape.line.fill.background()
        return shape

    today = date.today().strftime('%d %B %Y')
    total    = len(projects)
    active_c = sum(1 for p in projects if p.status.value == 'active')
    done_c   = sum(1 for p in projects if p.status.value == 'completed')
    hold_c   = sum(1 for p in projects if p.status.value == 'on_hold')
    budget   = sum(float(p.budget) for p in projects if p.budget)
    cmp      = _pdf_clean(company_name)

    # ── Slide 1: Cover ────────────────────────────────────────────────────
    s1 = slide()
    bg(s1, BRAND)
    rect(s1, 0, 2.9, 10, 0.04, RGBColor(0x38, 0xbd, 0xf8))
    tb(s1, 'AI PROJECT CONTROLS',  0.6, 0.5,  8.8, 0.5, size=9,  bold=False, color=SKY,   align='left')
    tb(s1, cmp,                     0.6, 1.0,  8.8, 1.5, size=30, bold=True,  align='left')
    tb(s1, 'Project Status Report', 0.6, 2.5,  8.8, 0.5, size=14, bold=False, color=SKY,  align='left')
    tb(s1, today,                   0.6, 3.15, 8.8, 0.4, size=10, bold=False, color=SLATE, align='left')

    # ── Slide 2: KPI Overview ────────────────────────────────────────────
    s2 = slide()
    bg(s2, RGBColor(0xf8, 0xfa, 0xfc))
    rect(s2, 0, 0, 10, 0.55, BRAND)
    tb(s2, 'Portfolio Overview', 0.4, 0.1, 9, 0.4, size=16, bold=True, align='left')

    kpis = [
        ('Total Projects', str(total),    BRAND),
        ('Active',         str(active_c), GREEN),
        ('Completed',      str(done_c),   RGBColor(0x47, 0x55, 0x69)),
        ('On Hold',        str(hold_c),   AMBER),
    ]
    for idx, (label, val, col) in enumerate(kpis):
        x = 0.4 + idx * 2.4
        rect(s2, x, 0.9, 2.15, 1.6, WHITE, RGBColor(0xd1, 0xd5, 0xdb))
        tb(s2, val,   x + 0.1, 1.0,  1.95, 0.8, size=30, bold=True,  color=col,   align='center')
        tb(s2, label, x + 0.1, 1.8,  1.95, 0.4, size=9,  bold=False, color=SLATE, align='center')

    if budget > 0:
        tb(s2, f'Total Portfolio Budget: Rs.{budget:,.0f}', 0.4, 2.75, 9, 0.45, size=13, bold=True, color=BRAND, align='center')

    # ── Slide 3: Projects Table ───────────────────────────────────────────
    s3 = slide()
    bg(s3, RGBColor(0xf8, 0xfa, 0xfc))
    rect(s3, 0, 0, 10, 0.55, BRAND)
    tb(s3, 'Project Status Summary', 0.4, 0.1, 9, 0.4, size=16, bold=True, align='left')

    if projects:
        visible = projects[:9]
        rows = len(visible) + 1
        tbl_shape = s3.shapes.add_table(rows, 5, Inches(0.3), Inches(0.75), Inches(9.4), Inches(4.0))
        tbl = tbl_shape.table

        col_widths = [3.0, 1.8, 1.5, 1.2, 1.9]
        for ci, w in enumerate(col_widths):
            tbl.columns[ci].width = Inches(w)

        headers = ['Project', 'Client', 'Status', 'Priority', 'Budget']
        for ci, h in enumerate(headers):
            cell = tbl.cell(0, ci)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = BRAND
            para = cell.text_frame.paragraphs[0]
            para.font.color.rgb = WHITE
            para.font.size = Pt(9)
            para.font.bold = True

        for ri, p in enumerate(visible, 1):
            vals = [
                p.name[:35],
                (p.client_name or '-')[:22],
                p.status.value.replace('_', ' ').title(),
                p.priority.value.title(),
                _fmt_budget(p.budget),
            ]
            for ci, val in enumerate(vals):
                cell = tbl.cell(ri, ci)
                cell.text = val
                para = cell.text_frame.paragraphs[0]
                para.font.size = Pt(8)
                para.font.color.rgb = RGBColor(0x1e, 0x29, 0x3b)
                if ri % 2 == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0xf1, 0xf5, 0xf9)

    # ── Slide 4: Closing ──────────────────────────────────────────────────
    s4 = slide()
    bg(s4, BRAND)
    rect(s4, 0, 2.5, 10, 0.04, RGBColor(0x38, 0xbd, 0xf8))
    tb(s4, cmp,                                     0.6, 1.0, 8.8, 1.0, size=26, bold=True,  align='center')
    tb(s4, 'Powered by AI Project Controls Platform', 0.6, 2.1, 8.8, 0.5, size=11, bold=False, color=SKY,   align='center')
    tb(s4, today,                                    0.6, 2.7, 8.8, 0.4, size=9,  bold=False, color=SLATE, align='center')

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf


# ── Email with PDF attachment ──────────────────────────────────────────────────

async def send_report_email(
    to_email: str,
    report_title: str,
    pdf_bytes: bytes,
    company_name: str,
    sender_name: str,
) -> bool:
    if not settings.SMTP_HOST or not settings.EMAILS_FROM_EMAIL:
        logger.warning('[EMAIL NOT SENT — No SMTP] Report email skipped.')
        return False

    safe_title = re.sub(r'[^\w\s\-]', '', report_title)
    filename = safe_title.replace(' ', '_') + '.pdf'

    msg = MIMEMultipart()
    msg['Subject'] = f'[{company_name}] {report_title}'
    msg['From']    = f'{settings.EMAILS_FROM_NAME} <{settings.EMAILS_FROM_EMAIL}>'
    msg['To']      = to_email

    html_body = f"""<!DOCTYPE html>
<html><head><meta charset="utf-8"></head>
<body style="font-family:Arial,sans-serif;background:#f4f4f4;margin:0;padding:0;">
  <div style="max-width:520px;margin:40px auto;background:#fff;border-radius:12px;overflow:hidden;box-shadow:0 2px 12px rgba(0,0,0,.1);">
    <div style="background:#1e3a5f;padding:28px 32px;">
      <h1 style="color:#fff;margin:0;font-size:20px;">AI Project Controls</h1>
      <p style="color:#bae6fd;margin:8px 0 0;font-size:13px;">{company_name}</p>
    </div>
    <div style="padding:32px;">
      <h2 style="color:#1e293b;margin:0 0 12px;">{report_title}</h2>
      <p style="color:#475569;">Hi,</p>
      <p style="color:#475569;">
        <strong>{sender_name}</strong> has shared a project report with you.
        Please find the PDF attached to this email.
      </p>
      <p style="color:#94a3b8;font-size:12px;margin-top:24px;">
        Sent from AI Project Controls Platform
      </p>
    </div>
  </div>
</body></html>"""

    msg.attach(MIMEText(html_body, 'html'))

    pdf_part = MIMEApplication(pdf_bytes, _subtype='pdf')
    pdf_part.add_header('Content-Disposition', 'attachment', filename=filename)
    msg.attach(pdf_part)

    try:
        await aiosmtplib.send(
            msg,
            hostname=settings.SMTP_HOST,
            port=settings.SMTP_PORT,
            username=settings.SMTP_USER,
            password=settings.SMTP_PASSWORD,
            start_tls=True,
        )
        return True
    except Exception as exc:
        logger.error(f'Failed to send report email to {to_email}: {exc}')
        return False
